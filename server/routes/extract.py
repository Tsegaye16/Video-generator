from fastapi import APIRouter, HTTPException
from models import ExtractRequest, ExtractionResponse, SlideData, ImageInfo, TableData
from config import settings, logger
from typing import List
import uuid
import matplotlib.pyplot as plt
from PIL import Image
import numpy as np
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from spire.presentation import Presentation as SpirePresentation, FileFormat
import logging
from wand.image import Image as WandImage
from PIL import Image as PILImage
import os
import io

def convert_wmf_to_png(wmf_path, png_path, resolution=300):
    """
    Convert WMF file to PNG using ImageMagick (wand)
    
    Args:
        wmf_path: Path to input WMF file
        png_path: Path to save output PNG
        resolution: DPI resolution for output (default 300)
    """
    try:
        # Convert using ImageMagick
        with WandImage(filename=wmf_path, resolution=resolution) as img:
            img.format = 'png'
            img.save(filename=png_path)
        
        # Optimize the PNG with Pillow
        with PILImage.open(png_path) as img:
            img.save(png_path, 'PNG', optimize=True)
        
        return True
    except Exception as e:
        print(f"Error converting WMF to PNG: {e}")
        return False

def process_wmf_image(blob, slide_number, image_index, output_dir):
    """
    Process a WMF image blob from PowerPoint
    
    Args:
        blob: Binary data of the WMF image
        slide_number: Slide number for naming
        image_index: Image index for naming
        output_dir: Directory to save files
    
    Returns:
        ImageInfo object or None if conversion fails
    """
    try:
        # Save original WMF
        wmf_filename = f"slide_{slide_number}_img_{image_index}.wmf"
        wmf_path = os.path.join(output_dir, wmf_filename)
        
        with open(wmf_path, 'wb') as f:
            f.write(blob)
        
        # Convert to PNG
        png_filename = f"slide_{slide_number}_img_{image_index}.png"
        png_path = os.path.join(output_dir, png_filename)
        
        if convert_wmf_to_png(wmf_path, png_path):
            # Get dimensions
            with PILImage.open(png_path) as img:
                width, height = img.size
            
            # Clean up WMF if PNG conversion succeeded
            os.remove(wmf_path)
            
            return ImageInfo(
                filename=png_filename,
                content_type="image/png",
                width_emu=width,
                height_emu=height
            )
        else:
            # Fallback to WMF if conversion fails
            with open(wmf_path, 'rb') as f:
                with PILImage.open(f) as img:
                    width, height = img.size
            
            return ImageInfo(
                filename=wmf_filename,
                content_type="image/wmf",
                width_emu=width,
                height_emu=height
            )
    
    except Exception as e:
        print(f"Error processing WMF image: {e}")
        return None
router = APIRouter()

import textwrap

def wrap_text(text: str, max_width: int = 12) -> str:
    return '\n'.join(textwrap.wrap(text, width=max_width))

# Configure logger
logger = logging.getLogger(__name__)

def process_group_shapes(shapes, slide_number, specific_extracted_path, image_index):
    """Recursively process shapes, including group shapes, to extract images."""
    images = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recursively process shapes within a group
                group_images = process_group_shapes(shape.shapes, slide_number, specific_extracted_path, image_index)
                images.extend(group_images)
                image_index += len(group_images)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, 'image'):
                try:
                    image = shape.image
                    if image and image.blob:
                        ext = image.ext.lower()
                        
                        # Special handling for WMF/EMF
                        if ext in ('wmf', 'emf'):
                            img_info = process_wmf_image(
                                image.blob,
                                slide_number,
                                image_index,
                                specific_extracted_path
                            )
                            if img_info:
                                images.append(img_info)
                                image_index += 1
                        else:
                            # Normal image processing for other formats
                            img_filename = f"slide_{slide_number}_img_{image_index}.{ext}"
                            img_path = os.path.join(specific_extracted_path, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(image.blob)
                            
                            # Convert to PNG if not already
                            if ext != 'png':
                                png_filename = f"slide_{slide_number}_img_{image_index}.png"
                                png_path = os.path.join(specific_extracted_path, png_filename)
                                
                                try:
                                    with PILImage.open(img_path) as img:
                                        img.save(png_path, 'PNG')
                                    os.remove(img_path)
                                    img_filename = png_filename
                                    ext = 'png'
                                except Exception as conv_e:
                                    print(f"Couldn't convert to PNG, keeping original: {conv_e}")
                            
                            with PILImage.open(os.path.join(specific_extracted_path, img_filename)) as img:
                                width, height = img.size
                            
                            images.append(ImageInfo(
                                filename=img_filename,
                                content_type=f"image/{ext}",
                                width_emu=width,
                                height_emu=height
                            ))
                            image_index += 1
                
                except Exception as e:
                    logger.error(f"Error extracting image from shape in slide {slide_number}: {e}")
            else:
                # Check for images in placeholders or other shapes
                if hasattr(shape, 'element') and shape.element.tag.endswith('pic'):
                    try:
                        image = shape.image
                        if image and image.blob:
                            # Same conversion process as above
                            original_ext = image.ext.lower()
                            original_img_filename = f"slide_{slide_number}_img_{image_index}_original.{original_ext}"
                            original_img_path = os.path.join(specific_extracted_path, original_img_filename)
                            with open(original_img_path, 'wb') as f:
                                f.write(image.blob)
                            
                            png_filename = f"slide_{slide_number}_img_{image_index}.png"
                            png_path = os.path.join(specific_extracted_path, png_filename)
                            
                            try:
                                with Image.open(original_img_path) as img:
                                    if img.mode in ('RGBA', 'LA', 'P'):
                                        img = img.convert('RGBA')
                                    elif img.mode != 'RGB':
                                        img = img.convert('RGB')
                                    img.save(png_path, 'PNG', quality=95)
                                
                                with Image.open(png_path) as img:
                                    width, height = img.size
                                
                                images.append(ImageInfo(
                                    filename=png_filename,
                                    content_type="image/png",
                                    width_emu=width,
                                    height_emu=height
                                ))
                                logger.info(f"Extracted and converted placeholder image to PNG: {png_filename}")
                                
                                # os.remove(original_img_path)
                                
                            except Exception as conv_e:
                                logger.error(f"Error converting placeholder image to PNG: {conv_e}")
                                with Image.open(original_img_path) as img:
                                    width, height = img.size
                                
                                images.append(ImageInfo(
                                    filename=original_img_filename,
                                    content_type=image.content_type,
                                    width_emu=width,
                                    height_emu=height
                                ))
                            
                            image_index += 1
                    except Exception as e:
                        logger.error(f"Error extracting placeholder image in slide {slide_number}: {e}")
        except Exception as e:
            logger.error(f"Error processing shape in slide {slide_number}: {e}")
    return images

def get_cell_background_color(cell):
    """Extract the background color of a cell."""
    fill = cell.fill
    if fill.type == 1:  # Solid fill
        if fill.fore_color.type == MSO_COLOR_TYPE.RGB:
            return fill.fore_color.rgb
        elif fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
            theme_color = fill.fore_color.theme_color
            # Map theme color to RGB (approximation, may need theme color mapping)
            # Placeholder: Return white if theme color is not mapped
            return RGBColor(255, 255, 255)
    return None  # Default to no color (transparent)

def get_cell_text_color(cell):
    """Extract the text color of a cell, handling all MSO_COLOR_TYPE values."""
    try:
        if cell.text_frame and cell.text_frame.paragraphs:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font and run.font.color and run.font.color.type:
                        color = run.font.color
                        if color.type == MSO_COLOR_TYPE.RGB:
                            return color.rgb
                        elif color.type == MSO_COLOR_TYPE.SCHEME:
                            # Map theme color to RGB using theme color dictionary
                            theme_color = color.theme_color
                            # Example theme color mapping (adjust based on presentation theme)
                            theme_rgb_map = {
                                MSO_THEME_COLOR.DARK_1: RGBColor(0, 0, 0),          # Black
                                MSO_THEME_COLOR.LIGHT_1: RGBColor(255, 255, 255),  # White
                                MSO_THEME_COLOR.DARK_2: RGBColor(68, 84, 106),     # Dark gray
                                MSO_THEME_COLOR.LIGHT_2: RGBColor(238, 236, 225),  # Light gray
                                MSO_THEME_COLOR.ACCENT_1: RGBColor(31, 73, 125),   # Blue
                                MSO_THEME_COLOR.ACCENT_2: RGBColor(79, 129, 189),  # Light blue
                                MSO_THEME_COLOR.ACCENT_3: RGBColor(192, 80, 77),   # Red
                                MSO_THEME_COLOR.ACCENT_4: RGBColor(155, 187, 89),  # Green
                                MSO_THEME_COLOR.ACCENT_5: RGBColor(128, 100, 162), # Purple
                                MSO_THEME_COLOR.ACCENT_6: RGBColor(75, 172, 198),  # Cyan
                                MSO_THEME_COLOR.HYPERLINK: RGBColor(0, 0, 255),    # Blue
                                MSO_THEME_COLOR.FOLLOWED_HYPERLINK: RGBColor(128, 0, 128),  # Purple
                                MSO_THEME_COLOR.TEXT_1: RGBColor(0, 0, 0),         # Black
                                MSO_THEME_COLOR.BACKGROUND_1: RGBColor(255, 255, 255),  # White
                                MSO_THEME_COLOR.TEXT_2: RGBColor(68, 84, 106),     # Dark gray
                                MSO_THEME_COLOR.BACKGROUND_2: RGBColor(238, 236, 225),  # Light gray
                            }
                            return theme_rgb_map.get(theme_color, RGBColor(0, 0, 0))
                        elif color.type == MSO_COLOR_TYPE.HSL:
                            # Convert HSL to RGB (placeholder, requires conversion logic)
                            return RGBColor(0, 0, 0)  # Default to black
                        elif color.type == MSO_COLOR_TYPE.PRESET:
                            # Preset colors are less common; default to black or map if known
                            return RGBColor(0, 0, 0)  # Default to black
                        elif color.type == MSO_COLOR_TYPE.SCRGB:
                            # scRGB uses floating-point RGB; convert to standard RGB
                            rgb = color.rgb  # May need scaling if scRGB is out of range
                            return rgb if rgb else RGBColor(0, 0, 0)
                        elif color.type == MSO_COLOR_TYPE.SYSTEM:
                            # System colors depend on OS; default to black
                            return RGBColor(0, 0, 0)  # Default to black
                        elif color.type == MSO_COLOR_TYPE.AUTO:
                            # AUTO typically inherits from theme or default (black)
                            return RGBColor(0, 0, 0)  # Default to black
        # Fallback for empty cells or no color defined
        return RGBColor(0, 0, 0)  # Default to black
    except Exception:
        # Handle unexpected errors
        return RGBColor(0, 0, 0)  # Default to black

@router.post("/api/extract")
async def extract_content(request: ExtractRequest):
    file_id = request.file_id
    original_file_path = os.path.join(settings.TEMP_UPLOAD_DIR, file_id)

    if not os.path.exists(original_file_path):
        logger.warning(f"File not found: {file_id}")
        raise HTTPException(status_code=404, detail=f"File not found: {file_id}")

    file_ext = os.path.splitext(file_id)[1].lower()
    processing_file_path = original_file_path
    converted_file_path = None

    if file_ext == '.ppt':
        try:
            spire_pptx = SpirePresentation()
            spire_pptx.LoadFromFile(original_file_path)
            converted_filename = f"{uuid.uuid4()}.pptx"
            converted_file_path = os.path.join(settings.TEMP_UPLOAD_DIR, converted_filename)
            spire_pptx.SaveToFile(converted_file_path, FileFormat.Pptx2019)
            spire_pptx.Dispose()
            processing_file_path = converted_file_path
            logger.info(f"Converted {file_id} to PPTX: {converted_filename}")
        except Exception as e:
            if os.path.exists(original_file_path):
                os.remove(original_file_path)
            logger.error(f"Error converting PPT to PPTX for {file_id}: {e}")
            raise HTTPException(status_code=500, detail=f"PPT to PPTX conversion failed: {e}")

    file_base_name = os.path.splitext(file_id)[0]
    specific_extracted_path = os.path.join(settings.EXTRACTED_CONTENT_DIR, file_base_name)
    os.makedirs(specific_extracted_path, exist_ok=True)

    extracted_slides_data: List[SlideData] = []

    try:
        prs = Presentation(processing_file_path)
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            current_slide_data = SlideData(slide_number=slide_number)
            image_index, table_index = 0, 0
            title_shape = None

            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_shape = slide.shapes.title
                current_slide_data.title = title_shape.text.strip()

            # Extract images, including those in groups or placeholders
            current_slide_data.images = process_group_shapes(slide.shapes, slide_number, specific_extracted_path, image_index)

            for shape in slide.shapes:
                if shape.has_text_frame and shape != title_shape:
                    text = shape.text.strip()
                    if text:
                        current_slide_data.text_elements.append(text)

                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_data_list = [
                        [wrap_text(cell.text.strip()) for cell in row.cells]
                        for row in table.rows
                    ]

                    # Extract cell background and text colors
                    cell_bg_colors = [
                        [get_cell_background_color(cell) for cell in row.cells]
                        for row in table.rows
                    ]
                    cell_text_colors = [
                        [get_cell_text_color(cell) for cell in row.cells] 
                        for row in table.rows
]

                    # Extract horizontal and vertical alignments
                    cell_h_alignments = []
                    cell_v_alignments = []
                    for row in table.rows:
                        row_h_alignments = []
                        row_v_alignments = []
                        for cell in row.cells:
                            # Get horizontal alignment
                            h_align = 'center'  # Default
                            if cell.text_frame and cell.text_frame.paragraphs:
                                paragraph = cell.text_frame.paragraphs[0]
                                if paragraph.alignment == PP_ALIGN.LEFT:
                                    h_align = 'left'
                                elif paragraph.alignment == PP_ALIGN.RIGHT:
                                    h_align = 'right'
                                elif paragraph.alignment == PP_ALIGN.CENTER:
                                    h_align = 'center'

                            # Get vertical alignment
                            v_align = 'center'  # Default
                            # if cell.vertical_anchor == MSO_VERTICAL_ANCHOR.TOP:
                            #     v_align = 'top'
                            # elif cell.vertical_anchor == MSO_VERTICAL_ANCHOR.BOTTOM:
                            #     v_align = 'bottom'
                            # elif cell.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE:
                            #     v_align = 'center'

                            row_h_alignments.append(h_align)
                            row_v_alignments.append(v_align)
                        cell_h_alignments.append(row_h_alignments)
                        cell_v_alignments.append(row_v_alignments)

                    if table_data_list:
                        num_rows = len(table_data_list)
                        num_cols = len(table_data_list[0]) if table_data_list else 1
                        fig_width = min(8, num_cols * 1.5)
                        fig_height = min(4, num_rows * 1.5)
                        fig, ax = plt.subplots(figsize=(fig_width, fig_height))
                        ax.axis('off')
                        plt.box(False)

                        # Create the table with matplotlib
                        table_img = ax.table(
                            cellText=table_data_list,
                            cellLoc='center',  # This will be overridden per cell
                            loc='center',
                            colWidths=[1.0 / num_cols] * num_cols
                        )
                        table_img.auto_set_font_size(False)
                        table_img.set_fontsize(9)

                        # Apply formatting to each cell
                        for row in range(num_rows):
                            for col in range(num_cols):
                                cell = table_img[(row, col)]
                                # Set background color
                                bg_color = cell_bg_colors[row][col]
                                if bg_color:
                                    cell.set_facecolor((bg_color[0] / 255, bg_color[1] / 255, bg_color[2] / 255))
                                else:
                                    cell.set_facecolor('white')  # Default to white if no color

                                # Set text color and alignment
                                text_color = cell_text_colors[row][col]
                                h_align = cell_h_alignments[row][col]
                                v_align = cell_v_alignments[row][col]
                                if text_color:
                                    cell.set_text_props(
                                        ha=h_align,
                                        va=v_align,
                                        color=(text_color[0] / 255, text_color[1] / 255, text_color[2] / 255)
                                    )
                                else:
                                    cell.set_text_props(
                                        ha=h_align,
                                        va=v_align,
                                        color='black'
                                    )

                                # Set cell properties
                                cell.set_linewidth(0.5)
                                cell.set_height(0.1)
                                cell.set_edgecolor('black')
                                cell.PAD = 0.05

                        table_img.auto_set_column_width([i for i in range(num_cols)])
                        fig.tight_layout(pad=0.0)
                        fig.subplots_adjust(left=0.0, right=1.0, top=1.0, bottom=0.0)

                        # Save the table as an image
                        img_filename = f"slide_{slide_number}_table_{table_index}.png"
                        img_save_path = os.path.join(specific_extracted_path, img_filename)
                        fig.savefig(img_save_path, dpi=300, bbox_inches='tight', pad_inches=0.1, transparent=True)
                        plt.close(fig)

                        # Add image info
                        with Image.open(img_save_path) as img:
                            width, height = img.size

                        current_slide_data.images.append(ImageInfo(
                            filename=img_filename,
                            content_type="image/png",
                            width_emu=width,
                            height_emu=height
                        ))

                        table_index += 1


            extracted_slides_data.append(current_slide_data)
            logger.info(f"Extracted slide {slide_number} with title: {current_slide_data.title}")

    except Exception as e:
        logger.error(f"Error extracting content from {file_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Extraction failed: {e}")

    finally:
        try:
            if os.path.exists(original_file_path):
                os.remove(original_file_path)
                logger.info(f"Cleaned up original file: {original_file_path}")
            if converted_file_path and os.path.exists(converted_file_path):
                os.remove(converted_file_path)
                logger.info(f"Cleaned up converted file: {converted_file_path}")
        except Exception as e:
            logger.warning(f"Failed to clean up files: {e}")

    return ExtractionResponse(
        file_id=file_id,
        extracted_content_path=specific_extracted_path,
        slides=extracted_slides_data
    )
